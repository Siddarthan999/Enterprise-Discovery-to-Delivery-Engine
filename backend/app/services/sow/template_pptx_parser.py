from pptx import Presentation


def extract_pptx_outline(path: str) -> dict:
    """Lightweight metadata extraction for uploaded pptx/potx templates —
    used only to show the user what they uploaded (slide count, roughly how
    many text areas each slide has). This does NOT attempt to understand
    slide content/meaning; that happens later at export time in
    pptx_generator.py, against the actual generated proposal content.
    """
    prs = Presentation(path)
    sections = []
    for i, slide in enumerate(prs.slides, start=1):
        text_shapes = [s for s in slide.shapes if getattr(s, "has_text_frame", False)]
        sections.append(
            f"Slide {i} ({len(text_shapes)} text area{'s' if len(text_shapes) != 1 else ''})"
        )

    return {
        "sections": sections,
        "slide_count": len(prs.slides),
    }