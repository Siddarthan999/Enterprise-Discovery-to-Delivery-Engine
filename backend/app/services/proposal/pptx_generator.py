import io
import copy
import textwrap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

EMU_PER_INCH = 914400

_PROTECTED_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER.SLIDE_NUMBER,
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.HEADER,
}

MIN_CONTENT_WIDTH_IN = 1.3
MIN_CONTENT_HEIGHT_IN = 0.35
MIN_CONTENT_AREA_IN2 = 0.6
MAX_ASPECT_RATIO = 2.5
GRID_AREA_RATIO_THRESHOLD = 2.0

# Fallback sizes ONLY used if the template's own slide master doesn't
# define title/body styles at all (rare — most real decks do). Whenever
# the master defines real sizes, those are used instead — see
# _master_default_sizes().
FALLBACK_TITLE_FONT_PT = 28
FALLBACK_BODY_FONT_PT = 18

TITLE_FONT_FLOOR_PT = 16
FONT_FLOOR_PT = 10
LINE_SPACING_FACTOR = 1.25
PARA_SPACING_IN = 0.06
SHAPE_BOTTOM_MARGIN_IN = 0.3

# A block is treated as "divider-worthy" (routed to a colorful, single-
# statement layout instead of the plain content layout) when it has just
# one short paragraph and no bullet list — matches how these layouts are
# actually shaped (large centered title + one subtitle-sized line).
DIVIDER_MAX_BODY_LINES = 1
DIVIDER_MAX_CHARS = 480


# ---------------------------------------------------------------------------
# Step 1: structured proposal -> ordered content blocks (unchanged)
# ---------------------------------------------------------------------------
def _build_content_blocks(sp: dict) -> list[dict]:
    blocks = []
    blocks.append({"title": sp.get("proposal_title") or "Proposal", "body": []})

    if sp.get("executive_summary"):
        blocks.append({"title": "Executive Summary", "body": [sp["executive_summary"]]})

    landscape_body = []
    if sp.get("current_landscape_intro"):
        landscape_body.append(sp["current_landscape_intro"])
    landscape_body.extend(sp.get("current_landscape_points") or [])
    if landscape_body:
        blocks.append({"title": "Current Landscape", "body": landscape_body})

    heard_body = []
    if sp.get("what_weve_heard_intro"):
        heard_body.append(sp["what_weve_heard_intro"])
    heard_body.extend(sp.get("what_weve_heard_themes") or [])
    if heard_body:
        blocks.append({"title": "What We've Heard", "body": heard_body})

    if sp.get("strategic_goals"):
        blocks.append({"title": "Strategic Goals", "body": sp["strategic_goals"]})

    if sp.get("target_outcomes"):
        blocks.append({"title": "Target Outcomes", "body": sp["target_outcomes"]})

    if sp.get("solution_overview") or sp.get("solution_components"):
        body = [sp["solution_overview"]] if sp.get("solution_overview") else []
        for comp in sp.get("solution_components") or []:
            name = comp.get("name") or "Component"
            desc = comp.get("description") or ""
            body.append(f"{name}: {desc}" if desc else name)
        blocks.append({"title": "Proposed Solution", "body": body})

    for i, phase in enumerate(sp.get("approach_phases") or [], start=1):
        title = phase.get("title") or f"Phase {i}"
        body = []
        if phase.get("objective"):
            body.append(f"Objective: {phase['objective']}")
        if phase.get("narrative"):
            body.append(phase["narrative"])
        body.extend(phase.get("indicative_activities") or [])
        if body:
            blocks.append({"title": f"Approach \u2014 {title}", "body": body})

    if sp.get("change_management_narrative"):
        blocks.append({"title": "Change Management Approach", "body": [sp["change_management_narrative"]]})

    if sp.get("risk_items"):
        body = []
        for item in sp["risk_items"]:
            risk = item.get("risk") or ""
            mitigation = item.get("mitigation") or ""
            if risk:
                body.append(f"{risk} (Mitigation: {mitigation})" if mitigation else risk)
        if body:
            blocks.append({"title": "Risk Management", "body": body})

    if sp.get("in_scope"):
        blocks.append({"title": "In Scope", "body": sp["in_scope"]})
    if sp.get("out_of_scope"):
        blocks.append({"title": "Out of Scope", "body": sp["out_of_scope"]})

    if sp.get("deliverables"):
        blocks.append({"title": "Deliverables", "body": sp["deliverables"]})

    commercial_body = []
    if sp.get("commercial_narrative"):
        commercial_body.append(sp["commercial_narrative"])
    for opt in sp.get("commercial_options") or []:
        name = opt.get("name") or "Option"
        desc = opt.get("description") or ""
        commercial_body.append(f"{name}: {desc}" if desc else name)
    if commercial_body:
        blocks.append({"title": "Commercial Approach", "body": commercial_body})

    team_body = []
    for role in sp.get("team_roles") or []:
        r = role.get("role") or "Role"
        d = role.get("description") or ""
        team_body.append(f"{r}: {d}" if d else r)
    if team_body:
        blocks.append({"title": "Team", "body": team_body})

    if sp.get("why_this_approach"):
        blocks.append({"title": "Why This Approach", "body": [sp["why_this_approach"]]})

    if sp.get("next_steps"):
        blocks.append({"title": "Recommended Next Steps", "body": sp["next_steps"]})

    return blocks


def _is_divider_worthy(body_lines: list[str]) -> bool:
    if not body_lines or len(body_lines) > DIVIDER_MAX_BODY_LINES:
        return False
    return len(body_lines[0]) <= DIVIDER_MAX_CHARS


# ---------------------------------------------------------------------------
# Step 2: layout selection
# ---------------------------------------------------------------------------
def _has_picture_placeholder(layout) -> bool:
    for ph in layout.placeholders:
        try:
            if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                return True
        except Exception:
            continue
    return False


def _select_layouts(prs):
    """Returns (cover_layout, default_layout, divider_layout | None,
    closing_layout | None). Only ever selects from layouts with NO picture
    placeholder — an unfilled picture placeholder renders as a large gray
    "click to add picture" box."""
    safe_layouts = [l for l in prs.slide_layouts if not _has_picture_placeholder(l)]
    if not safe_layouts:
        safe_layouts = list(prs.slide_layouts)

    by_name = {l.name: l for l in safe_layouts}

    cover_layout = by_name.get("Title Slide")
    if cover_layout is None:
        cover_layout = next((l for n, l in by_name.items() if "cover page" in n.lower()), None)
    if cover_layout is None:
        cover_layout = safe_layouts[0]

    default_layout = (
        by_name.get("Title and Content")
        or by_name.get("Section Header")
        or safe_layouts[0]
    )

    # Prefer a plain/neutral divider over a heavily-branded dark/colored
    # one so text contrast stays predictable regardless of which specific
    # brand palette the template uses.
    divider_layout = None
    for preferred in ("Section Divider Neutral", "Section Divider Evergreen", "Section Divider Dark"):
        if preferred in by_name:
            divider_layout = by_name[preferred]
            break
    if divider_layout is None:
        divider_layout = next(
            (l for n, l in by_name.items() if "section divider" in n.lower() and "image" not in n.lower()),
            None,
        )

    closing_layout = None
    for n, l in by_name.items():
        if "thank you" in n.lower():
            closing_layout = l
            break

    return cover_layout, default_layout, divider_layout, closing_layout


# ---------------------------------------------------------------------------
# Step 3: read the template's OWN designed font sizes from the slide
# master, so our shrink-to-fit logic starts from the real brand
# typography instead of generic fallback numbers.
# ---------------------------------------------------------------------------
def _master_default_sizes(slide_master) -> dict:
    sizes = {"title": None, "body": None}
    try:
        tx_styles = slide_master.element.find(qn("p:txStyles"))
        if tx_styles is None:
            return sizes
        for key, tag in (("title", "p:titleStyle"), ("body", "p:bodyStyle")):
            node = tx_styles.find(qn(tag))
            if node is None:
                continue
            lvl1 = node.find(qn("a:lvl1pPr"))
            if lvl1 is None:
                continue
            def_rpr = lvl1.find(qn("a:defRPr"))
            if def_rpr is not None and def_rpr.get("sz"):
                sizes[key] = int(def_rpr.get("sz")) / 100  # sz is in hundredths of a point
    except Exception:
        pass
    return sizes


# ---------------------------------------------------------------------------
# Step 4: shape filtering — decorative/icon shapes vs. genuine content areas
# (unchanged from previous fix)
# ---------------------------------------------------------------------------
def _is_vertical_text(shape) -> bool:
    try:
        body_pr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        if body_pr is not None:
            vert = body_pr.get("vert")
            if vert and vert != "horz":
                return True
    except Exception:
        pass
    return False


def _is_decorative_shape(shape) -> bool:
    width_in = (shape.width or 0) / EMU_PER_INCH
    height_in = (shape.height or 0) / EMU_PER_INCH

    if width_in <= 0 or height_in <= 0:
        return True
    if width_in < MIN_CONTENT_WIDTH_IN or height_in < MIN_CONTENT_HEIGHT_IN:
        return True
    if (width_in * height_in) < MIN_CONTENT_AREA_IN2:
        return True
    if height_in > width_in * MAX_ASPECT_RATIO:
        return True
    if _is_vertical_text(shape):
        return True
    return False


def _classify_slide_shapes(slide):
    title_shape = None
    content_candidates = []
    protected = []
    decorative = []

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        ph_type = None
        ph_idx = None
        if shape.is_placeholder:
            try:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
            except Exception:
                pass

        if ph_type in _PROTECTED_PLACEHOLDER_TYPES:
            protected.append(shape)
            continue

        if title_shape is None and (
            ph_idx == 0 or ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
        ):
            title_shape = shape
            continue

        if _is_decorative_shape(shape):
            decorative.append(shape)
            continue

        content_candidates.append(shape)

    if title_shape is None and content_candidates:
        content_candidates.sort(key=lambda s: (s.top if s.top is not None else 0))
        title_shape = content_candidates.pop(0)

    content_candidates.sort(
        key=lambda s: (s.top if s.top is not None else 0, s.left if s.left is not None else 0)
    )

    return title_shape, content_candidates, protected, decorative


def _add_fallback_textboxes(slide, prs, title_default_pt, body_default_pt):
    left = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, Inches(0.4), prs.slide_width - Inches(1.2), Inches(1))
    title_box.text_frame.paragraphs[0].font.size = Pt(title_default_pt)
    title_box.text_frame.paragraphs[0].font.bold = True

    body_box = slide.shapes.add_textbox(left, Inches(1.6), prs.slide_width - Inches(1.2), prs.slide_height - Inches(2.2))
    body_box.text_frame.word_wrap = True
    body_box.text_frame.paragraphs[0].font.size = Pt(body_default_pt)

    return title_box, [body_box]


# ---------------------------------------------------------------------------
# Step 5: write text while preserving styling
# ---------------------------------------------------------------------------
def _shape_has_text(text_frame) -> bool:
    return any(r.text.strip() for p in text_frame.paragraphs for r in p.runs)


def _set_text_preserving_style(text_frame, text: str):
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        text_frame.text = text
        return
    p0 = paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in p0.runs[1:]:
            r.text = ""
    else:
        run = p0.add_run()
        run.text = text
    for p in paragraphs[1:]:
        for r in p.runs:
            r.text = ""


def _set_bullets_preserving_style(text_frame, lines: list[str]):
    lines = [l for l in lines if l and l.strip()] or [""]
    txBody = text_frame._txBody
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        text_frame.text = lines[0]
        paragraphs = list(text_frame.paragraphs)

    template_p_elem = paragraphs[0]._p
    for p in paragraphs[1:]:
        txBody.remove(p._p)

    first_para = text_frame.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = lines[0]
        for r in first_para.runs[1:]:
            r.text = ""
    else:
        run = first_para.add_run()
        run.text = lines[0]

    for line in lines[1:]:
        new_elem = copy.deepcopy(template_p_elem)
        txBody.append(new_elem)
        new_para = text_frame.paragraphs[-1]
        if new_para.runs:
            new_para.runs[0].text = line
            for r in new_para.runs[1:]:
                r.text = ""
        else:
            run = new_para.add_run()
            run.text = line


def _clear_text_frame(text_frame):
    if not _shape_has_text(text_frame):
        return  # already empty — leave the native "click to add text" prompt intact
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        return
    for p in paragraphs[1:]:
        text_frame._txBody.remove(p._p)
    p0 = text_frame.paragraphs[0]
    for r in p0.runs:
        r.text = ""
    if not p0.runs:
        p0.add_run().text = ""


# ---------------------------------------------------------------------------
# NEW — Step 5b: strip any placeholder shape that ends up with no text.
#
# Any placeholder left empty (a decorative "eyebrow" placeholder we never
# had content for, an unused Subtitle on the cover slide, an extra grid
# cell with no matching bullet, etc.) still shows as a native "Click to
# add text" prompt the moment someone opens the deck in PowerPoint's
# Normal view — LibreOffice/PDF export hides it, which is why it wasn't
# visible in a rendered preview but is very visible in the real file.
# We remove those shapes outright instead of leaving them empty.
# Slide number / date / footer / header placeholders are left alone
# since their visible text comes from field codes, not runs.
# ---------------------------------------------------------------------------
def _remove_empty_placeholders(slide):
    to_remove = []
    for shape in slide.shapes:
        if not shape.is_placeholder or not getattr(shape, "has_text_frame", False):
            continue
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            ph_type = None
        if ph_type in _PROTECTED_PLACEHOLDER_TYPES:
            continue
        if _shape_has_text(shape.text_frame):
            continue
        to_remove.append(shape)
    for shape in to_remove:
        shape._element.getparent().remove(shape._element)


# ---------------------------------------------------------------------------
# NEW — Step 5c: safely grow a shape's height without corrupting its
# geometry.
#
# ROOT CAUSE OF THE "STRETCHED INTO A VERTICAL LINE" BUG:
#
# Many placeholder shapes (e.g. a Subtitle placeholder on a title-style
# layout) have NO local <a:xfrm> of their own in the slide XML — they
# simply inherit left/top/width/height from the layout's placeholder.
# python-pptx's read-side `shape.width` / `shape.height` correctly
# resolve that inheritance and return the real values.
#
# But python-pptx's *write*-side `shape.height = ...` setter does this
# under the hood (see pptx.oxml.shapes.shared.BaseShapeElement.cy):
#     xfrm = self.get_or_add_xfrm()   # creates a BRAND NEW <a:xfrm>
#     xfrm.cy = value                 # only cy gets set
# A freshly created <a:xfrm> defaults every other attribute (cx, x, y)
# to 0. Because a placeholder shape treats "locally set to 0" as a real,
# directly-applied value (not "unset, fall back to the layout"), the
# shape's effective WIDTH permanently collapses to 0 the instant we set
# only its height — producing exactly the "text box squeezed into a
# vertical line" artifact seen on the Executive Summary / Why This
# Approach slides.
#
# This isn't specific to this one template — it will happen on ANY
# template whenever a placeholder without its own local xfrm needs its
# height grown to fit body text. The fix is to snapshot the shape's
# effective left/top/width BEFORE touching height, then re-apply them
# immediately after, so the newly-created xfrm ends up fully populated
# instead of half-zeroed.
# ---------------------------------------------------------------------------
def _grow_shape_height(shape, new_height_emu: int):
    # Snapshot the shape's CURRENT effective geometry first. For
    # placeholders this resolves inheritance from the layout; for
    # regular shapes it's just whatever is already set.
    left = shape.left
    top = shape.top
    width = shape.width

    shape.height = Emu(int(new_height_emu))

    # Always re-write left/top/width explicitly — not only when
    # python-pptx's own getters claim they changed. python-pptx's read
    # side is happy to fall back to the layout for any xfrm sub-element
    # (<a:off> or <a:ext>) that's simply absent, but that per-attribute
    # inheritance isn't reliably honored by every renderer once a shape
    # has ANY local <a:xfrm> — PowerPoint in particular has a long-
    # standing inconsistency here. A shape with a local <a:ext> but no
    # <a:off> can render pinned near (0,0) in PowerPoint even though
    # LibreOffice/python-pptx correctly infer the missing offset from
    # the layout. Writing all four values unconditionally makes the
    # xfrm fully self-contained so there's no ambiguity left for any
    # renderer to resolve differently.
    if width is not None:
        shape.width = width
    if left is not None:
        shape.left = left
    if top is not None:
        shape.top = top


# ---------------------------------------------------------------------------
# Step 6: fit text to its shape — anchored to the TEMPLATE's own
# designed font size, not a generic constant.
# ---------------------------------------------------------------------------
def _get_base_font_size(text_frame, default_pt) -> float:
    for p in text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                return r.font.size.pt
    return default_pt


def _estimate_block_height_in(lines: list[str], shape_width_in: float, font_pt: float) -> float:
    avg_char_width_in = (font_pt * 0.52) / 72.0
    usable_width_in = max(shape_width_in - 0.2, 0.5)
    chars_per_line = max(int(usable_width_in / avg_char_width_in), 8)

    line_height_in = (font_pt * LINE_SPACING_FACTOR) / 72.0
    total_lines = 0
    for line in lines:
        wrapped = textwrap.wrap(line, width=chars_per_line) or [""]
        total_lines += len(wrapped)

    return total_lines * line_height_in + max(len(lines) - 1, 0) * PARA_SPACING_IN


def _fit_text_to_shape(shape, slide, prs, lines: list[str], font_floor_pt: float, default_font_pt: float):
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    shape_width_in = (shape.width or Emu(0)) / EMU_PER_INCH
    shape_height_in = (shape.height or Emu(0)) / EMU_PER_INCH
    base_font_pt = _get_base_font_size(text_frame, default_font_pt)

    # Walk down from the template's actual designed size in whole points,
    # only shrinking as far as genuinely needed to fit.
    chosen_pt = font_floor_pt
    candidate = int(round(base_font_pt))
    floor = int(round(font_floor_pt))
    found = False
    for candidate_pt in range(candidate, floor - 1, -1):
        needed_in = _estimate_block_height_in(lines, shape_width_in, candidate_pt)
        if needed_in <= shape_height_in:
            chosen_pt = candidate_pt
            found = True
            break
    if not found:
        chosen_pt = floor

    needed_at_floor_in = _estimate_block_height_in(lines, shape_width_in, floor)
    if needed_at_floor_in > shape_height_in:
        slide_height_in = prs.slide_height / EMU_PER_INCH
        shape_top_in = (shape.top or Emu(0)) / EMU_PER_INCH
        max_available_in = slide_height_in - shape_top_in - SHAPE_BOTTOM_MARGIN_IN
        new_height_in = min(needed_at_floor_in, max_available_in)
        if new_height_in > shape_height_in:
            # NOTE: use _grow_shape_height(), not `shape.height = ...`
            # directly — see the docstring on _grow_shape_height() for
            # why a plain assignment silently zeroes the shape's width
            # on placeholders that inherit their geometry from the layout.
            _grow_shape_height(shape, new_height_in * EMU_PER_INCH)
        chosen_pt = floor

    for p in text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(chosen_pt)

    try:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Step 7: fill a single slide
# ---------------------------------------------------------------------------
def _fill_slide(slide, prs, title: str, body_lines: list[str]):
    master_sizes = _master_default_sizes(slide.slide_layout.slide_master)
    title_default_pt = master_sizes.get("title") or FALLBACK_TITLE_FONT_PT
    body_default_pt = master_sizes.get("body") or FALLBACK_BODY_FONT_PT

    title_shape, content_shapes, _protected, _decorative = _classify_slide_shapes(slide)

    if title_shape is None and not content_shapes:
        title_shape, content_shapes = _add_fallback_textboxes(slide, prs, title_default_pt, body_default_pt)

    if title_shape is not None:
        _set_text_preserving_style(title_shape.text_frame, title)
        _fit_text_to_shape(title_shape, slide, prs, [title], TITLE_FONT_FLOOR_PT, title_default_pt)

    if not content_shapes and body_lines:
        _, fallback = _add_fallback_textboxes(slide, prs, title_default_pt, body_default_pt)
        content_shapes = fallback

    if not body_lines:
        for shape in content_shapes:
            _clear_text_frame(shape.text_frame)
        _remove_empty_placeholders(slide)
        return

    if len(content_shapes) == 1:
        target = content_shapes[0]
        _set_bullets_preserving_style(target.text_frame, body_lines)
        _fit_text_to_shape(target, slide, prs, body_lines, FONT_FLOOR_PT, body_default_pt)
        _remove_empty_placeholders(slide)
        return

    areas = [max((s.width or 1) * (s.height or 1), 1) for s in content_shapes]
    is_grid = (max(areas) / min(areas)) < GRID_AREA_RATIO_THRESHOLD

    if is_grid and len(body_lines) >= 2:
        n = len(content_shapes)
        for i, shape in enumerate(content_shapes):
            if i >= len(body_lines):
                _clear_text_frame(shape.text_frame)
                continue
            item_lines = body_lines[i:] if (i == n - 1 and len(body_lines) > n) else [body_lines[i]]
            _set_bullets_preserving_style(shape.text_frame, item_lines)
            _fit_text_to_shape(shape, slide, prs, item_lines, FONT_FLOOR_PT, body_default_pt)
    else:
        largest = max(content_shapes, key=lambda s: (s.width or 0) * (s.height or 0))
        _set_bullets_preserving_style(largest.text_frame, body_lines)
        _fit_text_to_shape(largest, slide, prs, body_lines, FONT_FLOOR_PT, body_default_pt)
        for shape in content_shapes:
            if shape is not largest:
                _clear_text_frame(shape.text_frame)

    # Any placeholder shape (title, unused grid cell, decorative "eyebrow"
    # tag, etc.) that still has no text at this point gets removed rather
    # than left as a native empty prompt.
    _remove_empty_placeholders(slide)


# ---------------------------------------------------------------------------
# Step 8: remove the template's original slides once ours have been added.
#
# FIX: the previous version only stripped the <p:sldId> entries from
# presentation.xml's slide list. It never dropped the corresponding
# relationship in presentation.xml.rels, so the original template's slide
# parts (and everything they reference — images, layouts rels, etc.)
# stayed physically present in the .pptx package forever, just no longer
# shown. For a ~14-slide deck built from a large template this silently
# left ~70 orphaned slide parts sitting in the file:
#   - it bloats the file size for no reason (5MB -> 3MB after the fix
#     for this specific file)
#   - reopening/resaving the file with python-pptx (or some other tools)
#     can produce colliding/duplicate part names once new slides are
#     added later, which is the kind of inconsistency that makes
#     PowerPoint flag the file as needing repair.
# Dropping the relationship via `part.drop_rel()` lets python-pptx's
# package writer naturally exclude the now-unreferenced part on save —
# no manual part surgery required.
# ---------------------------------------------------------------------------
def _delete_first_n_slides(prs, n: int):
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    part = prs.part
    for sld in slide_ids[:n]:
        rId = sld.get(qn("r:id"))
        xml_slides.remove(sld)
        if rId:
            part.drop_rel(rId)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------
def generate_proposal_pptx(template_path: str, structured_proposal: dict) -> bytes:
    prs = Presentation(template_path)
    blocks = _build_content_blocks(structured_proposal)

    original_slide_count = len(prs.slides)
    if original_slide_count == 0:
        raise ValueError("The selected template has no slides")

    cover_layout, default_layout, divider_layout, closing_layout = _select_layouts(prs)

    plan = []
    for i, block in enumerate(blocks):
        if i == 0:
            layout = cover_layout
        elif divider_layout is not None and _is_divider_worthy(block.get("body") or []):
            layout = divider_layout
        else:
            layout = default_layout
        plan.append((layout, block["title"], block.get("body", [])))

    if closing_layout is not None:
        plan.append((
            closing_layout,
            "Thank You",
            ["We look forward to the opportunity to work together."],
        ))

    new_slides = [prs.slides.add_slide(layout) for layout, _t, _b in plan]

    for slide, (_layout, title, body) in zip(new_slides, plan):
        _fill_slide(slide, prs, title, body)

    _delete_first_n_slides(prs, original_slide_count)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()