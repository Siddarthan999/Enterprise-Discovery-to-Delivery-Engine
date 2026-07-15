from pypdf import PdfReader
import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import os
import re
from email import policy
from email.parser import BytesParser

# ---------------- CLEAN TEXT ----------------
def clean_text(text: str) -> str:
    if not text:
        return ""
    return text.replace("\x00", "").strip()


# ---------------- TXT ----------------
def parse_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()

    return {
        "title": os.path.basename(file_path),
        "content": clean_text(text),
        "type": "txt"
    }

# ---------------- PDF ----------------
def parse_pdf(file_path):
    try:
        reader = PdfReader(file_path)
    except Exception as e:
        return {
            "title": os.path.basename(file_path),
            "content": f"PDF_READ_FAILED: {e}",
            "type": "pdf"
        }

    if reader.is_encrypted:
        # Try an empty password first — many "encrypted" PDFs are just
        # permission-restricted with no real password, which pypdf still
        # flags as encrypted.
        try:
            reader.decrypt("")
        except Exception:
            return {
                "title": os.path.basename(file_path),
                "content": "PDF_ENCRYPTED_COULD_NOT_DECRYPT",
                "type": "pdf"
            }

    text_pages = []

    for page in reader.pages:
        try:
            extracted = page.extract_text()
            if extracted:
                text_pages.append(extracted)
        except Exception:
            continue

    text = "\n".join(text_pages)

    # fallback guard
    if len(text.strip()) < 50:
        text = "PDF_TEXT_EXTRACTION_FAILED_OR_SCANNED_DOCUMENT"

    return {
        "title": os.path.basename(file_path),
        "content": clean_text(text),
        "type": "pdf"
    }

# ---------------- DOCX ----------------
def _extract_docx_tables(doc) -> list:
    """python-docx's .paragraphs only covers body paragraphs — tables are
    a completely separate collection and were never being read at all."""
    table_texts = []
    for table in doc.tables:
        rows_text = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows_text.append(" | ".join(cells))
        if rows_text:
            table_texts.append("\n".join(rows_text))
    return table_texts


def parse_docx(file_path):
    doc = docx.Document(file_path)
    paragraph_text = "\n".join([p.text for p in doc.paragraphs])

    table_blocks = _extract_docx_tables(doc)
    table_text = "\n\n--- Table ---\n".join(table_blocks) if table_blocks else ""

    combined = paragraph_text
    if table_text:
        combined += "\n\n--- Table ---\n" + table_text

    return {
        "title": os.path.basename(file_path),
        "content": clean_text(combined),
        "type": "docx"
    }


# ---------------- PPTX ----------------
def _extract_shape_text(shape) -> list:
    """hasattr(shape, 'text') alone misses grouped shapes (their text lives
    on the child shapes, not the group itself) and slide tables (which use
    a different API than text frames)."""
    texts = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            texts.extend(_extract_shape_text(sub_shape))
        return texts

    if shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            texts.append(" | ".join(cells))
        return texts

    if shape.has_text_frame and shape.text_frame.text:
        texts.append(shape.text_frame.text)

    return texts


def parse_pptx(file_path):
    prs = Presentation(file_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            text.extend(_extract_shape_text(shape))

        # Speaker notes often carry the actual explanation behind a slide's
        # bullet points and were never being read before.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                text.append(f"[Speaker notes] {notes.strip()}")

    return {
        "title": os.path.basename(file_path),
        "content": clean_text("\n".join(text)),
        "type": "pptx"
    }


# ---------------- CSV ----------------
def parse_csv(file_path):
    df = pd.read_csv(file_path)
    text = df.astype(str).to_string(index=False)

    return {
        "title": os.path.basename(file_path),
        "content": clean_text(text),
        "type": "csv"
    }


# ---------------- EXCEL ----------------
def parse_excel(file_path):
    # pd.read_excel() with no sheet_name only reads the FIRST sheet —
    # sheet_name=None returns a dict of every sheet in the workbook.
    sheets = pd.read_excel(file_path, sheet_name=None)

    blocks = []
    for sheet_name, df in sheets.items():
        sheet_text = df.astype(str).to_string(index=False)
        blocks.append(f"--- Sheet: {sheet_name} ---\n{sheet_text}")

    text = "\n\n".join(blocks)

    return {
        "title": os.path.basename(file_path),
        "content": clean_text(text),
        "type": "excel"
    }


# ---------------- MARKDOWN ----------------
def parse_md(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()

    return {
        "title": os.path.basename(file_path),
        "content": clean_text(text),
        "type": "md"
    }


# ---------------- EMAIL ----------------
def _strip_html_tags(html: str) -> str:
    """Minimal HTML-to-text fallback for HTML-only email bodies."""
    text = re.sub(r"<(script|style)[^>]*>.*?</\1>", "", html, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def parse_eml(file_path):
    # Must open in binary mode — BytesParser needs raw bytes to correctly
    # detect and decode Content-Transfer-Encoding (base64, quoted-printable).
    # Reading as plain text (the old behavior) is exactly why base64 MIME
    # bodies were ending up verbatim in the extracted content.
    with open(file_path, "rb") as f:
        msg = BytesParser(policy=policy.default).parsebytes(f.read())

    parts = []

    subject = msg.get("subject", "")
    sender = msg.get("from", "")
    if subject:
        parts.append(f"Subject: {subject}")
    if sender:
        parts.append(f"From: {sender}")

    if msg.is_multipart():
        walked = list(msg.walk())
        has_plain = any(p.get_content_type() == "text/plain" for p in walked)

        for part in walked:
            content_type = part.get_content_type()
            content_disposition = str(part.get("Content-Disposition") or "")

            # Skip attachments — we only want inline readable text.
            if "attachment" in content_disposition:
                continue

            if content_type == "text/plain":
                payload = part.get_payload(decode=True)  # auto-decodes base64/QP
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    try:
                        parts.append(payload.decode(charset, errors="replace"))
                    except LookupError:
                        parts.append(payload.decode("utf-8", errors="replace"))

            elif content_type == "text/html" and not has_plain:
                # Only fall back to HTML if there's no plain-text part at all.
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    html = payload.decode(charset, errors="replace")
                    parts.append(_strip_html_tags(html))
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            charset = msg.get_content_charset() or "utf-8"
            parts.append(payload.decode(charset, errors="replace"))

    text = "\n\n".join(p.strip() for p in parts if p and p.strip())

    return {
        "title": os.path.basename(file_path),
        "content": clean_text(text),
        "type": "eml"
    }


# ---------------- .vtt (CAPTIONS) ----------------
def parse_vtt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()

    # Remove WEBVTT header + timestamps
    cleaned_lines = []

    for line in raw.split("\n"):
        line = line.strip()

        # skip headers / timestamps / empty lines
        if (
            line.startswith("WEBVTT")
            or "-->" in line
            or re.match(r"^\d+$", line)
            or line == ""
        ):
            continue

        # remove speaker tags like <v Name>
        line = re.sub(r"<v .*?>", "", line)
        line = line.replace("</v>", "")

        cleaned_lines.append(line)

    return {
        "title": os.path.basename(file_path),
        "content": " ".join(cleaned_lines),
        "type": "vtt"
    }

# ---------------- ROUTER ----------------
def parse_document(file_path):
    ext = file_path.split(".")[-1].lower()

    type_map = {
        "txt": "txt",
        "pdf": "pdf",
        "docx": "docx",
        "pptx": "pptx",
        "csv": "csv",
        "xls": "excel",
        "xlsx": "excel",
        "md": "md",
        "eml": "eml",
        "vtt": "vtt"
    }

    if ext == "txt":
        doc = parse_txt(file_path)
    elif ext == "pdf":
        doc = parse_pdf(file_path)
    elif ext == "docx":
        doc = parse_docx(file_path)
    elif ext == "pptx":
        doc = parse_pptx(file_path)
    elif ext == "csv":
        doc = parse_csv(file_path)
    elif ext in ["xls", "xlsx"]:
        doc = parse_excel(file_path)
    elif ext == "md":
        doc = parse_md(file_path)
    elif ext == "eml":
        doc = parse_eml(file_path)
    elif ext == "vtt":
        doc = parse_vtt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    doc["type"] = type_map.get(ext, ext)
    return doc

# You now have a unified enterprise ingestion parser:
"""
PDF     → text
DOCX    → text
PPTX    → slides → text
CSV     → table → text
Excel   → sheet → text
MD      → raw text
EML     → email content (properly MIME-decoded)
"""